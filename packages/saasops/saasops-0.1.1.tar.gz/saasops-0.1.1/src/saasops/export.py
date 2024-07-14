import pandas as pd

import saasops.calc as calc
import saasops.database as database
import saasops.visualization as visualization

from pptx import Presentation
from pptx.util import Inches

# Extract and store database data in Excel workbook


def export_data_to_xlsx(engine, start_date, end_date):
    df_customers = database.fetch_data_from_db(engine, "customers")
    df_contracts = database.fetch_data_from_db(engine, "contracts")
    df_segments = database.fetch_data_from_db(engine, "segments")
    df_invoicesegments = database.fetch_data_from_db(engine, "invoicesegments")
    df_invoices = database.fetch_data_from_db(engine, "invoices")

    print("Exporting data to Excel workbook...")
    print(f"  - Customers: {len(df_customers)} rows")
    print(f"  - Contracts: {len(df_contracts)} rows")
    print(f"  - Segments: {len(df_segments)} rows")
    print(f"  - Invoices: {len(df_invoices)} rows")
    print(f"  - Invoice Segments: {len(df_invoicesegments)} rows")

    df_invoices = pd.merge(
        df_invoices,
        df_invoicesegments[["invoiceid", "segmentid"]],
        on="invoiceid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices,
        df_segments[["segmentid", "contractid"]],
        on="segmentid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices,
        df_contracts[["contractid", "customerid"]],
        on="contractid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices, df_customers[["customerid", "name"]], on="customerid", how="left"
    )
    columns_to_drop = ["customerid", "segmentid", "contractid"]
    df_invoices.drop(columns=columns_to_drop, inplace=True)

    df_contracts = pd.merge(
        df_contracts, df_customers[["customerid", "name"]], on="customerid", how="left"
    )
    df_contracts.drop("customerid", axis=1, inplace=True)

    df_revenue = calc.populate_revenue_df(start_date, end_date, "mid", engine)
    df_metrics = calc.populate_metrics_df(start_date, end_date, engine)
    df_carrarr = calc.populate_bkings_carr_arr_df(start_date, end_date, engine)

    df_arr = calc.customer_arr_df(end_date, engine)
    df_carr = calc.customer_carr_df(end_date, engine)

    # Create a Pandas Excel writer using XlsxWriter as the engine.
    with pd.ExcelWriter("exports/data.xlsx", engine="xlsxwriter") as writer:
        df_customers.to_excel(writer, sheet_name="customers", index=False)
        df_contracts.to_excel(writer, sheet_name="contracts", index=False)
        df_invoices.to_excel(writer, sheet_name="invoices", index=False)
        df_revenue.to_excel(writer, sheet_name="revenue", index=True)
        df_metrics.to_excel(writer, sheet_name="metrics", index=True)
        df_carrarr.to_excel(writer, sheet_name="carrarr", index=True)
        df_arr.to_excel(writer, sheet_name="arr", index=True)
        df_carr.to_excel(writer, sheet_name="carr", index=True)


def export_data_to_pptx(engine, start_date, end_date):
    df_customers = database.fetch_data_from_db(engine, "customers")
    df_contracts = database.fetch_data_from_db(engine, "contracts")
    df_segments = database.fetch_data_from_db(engine, "segments")
    df_invoicesegments = database.fetch_data_from_db(engine, "invoicesegments")
    df_invoices = database.fetch_data_from_db(engine, "invoices")

    print("Exporting data to PowerPoint presentation...")
    print(f"  - Customers: {len(df_customers)} rows")
    print(f"  - Contracts: {len(df_contracts)} rows")
    print(f"  - Segments: {len(df_segments)} rows")
    print(f"  - Invoices: {len(df_invoices)} rows")
    print(f"  - Invoice Segments: {len(df_invoicesegments)} rows")

    df_invoices = pd.merge(
        df_invoices,
        df_invoicesegments[["invoiceid", "segmentid"]],
        on="invoiceid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices,
        df_segments[["segmentid", "contractid"]],
        on="segmentid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices,
        df_contracts[["contractid", "customerid"]],
        on="contractid",
        how="left",
    )
    df_invoices = pd.merge(
        df_invoices, df_customers[["customerid", "name"]], on="customerid", how="left"
    )
    columns_to_drop = ["customerid", "segmentid", "contractid"]
    df_invoices.drop(columns=columns_to_drop, inplace=True)

    df_contracts = pd.merge(
        df_contracts, df_customers[["customerid", "name"]], on="customerid", how="left"
    )
    df_contracts.drop("customerid", axis=1, inplace=True)

    visualization.create_mrr_change_chart(engine, start_date, end_date)
    visualization.create_monthly_mrr_chart(engine, start_date, end_date)
    visualization.ttm_ndr_gdr_chart(engine, end_date)
    visualization.create_bookings_arr_carr_chart(engine, start_date, end_date)

    # Create a new presentation
    prs = Presentation("exports/ppt_template.pptx")
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    chart_slides_data = [
        {"title": "Bookings, ARR, CARR", "image_path": "exports/bookings_arr_carr.png"},
        {"title": "Monthly MRR", "image_path": "exports/monthly_mrr.png"},
        {"title": "MRR Change", "image_path": "exports/mrr_change.png"},
        {
            "title": "TTM NDR & GDR",
            "image_path": "exports/trailing_12_month_values.png",
        },
    ]

    for slide_data in chart_slides_data:
        add_chart_slide(prs, slide_data["title"], slide_data["image_path"])

    prs.save("exports/data_export.pptx")


def export_chart_images(
    engine,
    start_date,
    end_date,
    customer=None,
    contract=None,
    show_gridlines=False,
    frequency="M",
    ignoreoverrides=False,
):
    # Generate charts
    visualization.create_mrr_change_chart(
        engine, start_date, end_date, customer, contract, frequency
    )
    visualization.create_monthly_mrr_chart(
        engine, start_date, end_date, customer, contract, show_gridlines, frequency
    )
    visualization.ttm_ndr_gdr_chart(engine, end_date, customer, contract)
    visualization.create_bookings_arr_carr_chart(
        engine,
        start_date,
        end_date,
        customer,
        contract,
        show_gridlines,
        frequency,
        ignoreoverrides,
    )
    visualization.create_monthly_arr_chart(
        engine,
        start_date,
        end_date,
        customer,
        contract,
        show_gridlines,
        frequency,
        ignoreoverrides,
    )


def add_chart_slide(prs, slide_title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[23])
    title = slide.shapes.title
    title.text = slide_title
    # left = Inches(0.5)
    # top = Inches(1.5)
    # height = Inches(5)
    # pic = slide.shapes.add_picture(image_path, left, top, height=height) UNUSED??
