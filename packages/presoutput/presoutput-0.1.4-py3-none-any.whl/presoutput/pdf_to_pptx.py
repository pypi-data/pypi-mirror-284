from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches
import io


def pdf_to_pptx(
    input_pdf_path: str,
    output_pptx_path: str,
    dpi: int = 300,
    pptx_height: float = 7.5,
    pptx_width: float = 13.33,
    verbose: int = 0,
):
    """
    Convert a PDF file to a PowerPoint presentation.

    Args:
        input_pdf_path (str): Path to the input PDF file.
        output_pptx_path (str): Path to save the output PowerPoint presentation.
        dpi (int, optional): DPI (dots per inch) for image conversion. Default is 300.
        pptx_height(float, optional): PPTX height in inches. Default is 7.5 (standard widescreen presentation size)
        pptx_width(float, optional): PPTX width in inches. Default is 13.33 (standard widescreen presentation size)
        verbose (bool, optional): If True, show progress messages. Default is True.

    Retuns:
        None

    Examples:
        >>> pdf_to_pptx('document.pdf', 'output.pptx', dpi=500)
    """

    # Convert PDF pages to images
    images = convert_from_path(input_pdf_path, dpi=dpi)

    # Create a PowerPoint presentation object
    presentation = Presentation()

    # Set slide dimensions to 16:9
    presentation.slide_width = Inches(pptx_width)
    presentation.slide_height = Inches(pptx_height)

    if verbose >= 1:
        print(f"Converting PDF to PPTX. Total pages: {len(images)}")

    for i, image in enumerate(images, start=1):
        # Save image to a byte array
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format="PNG")
        img_byte_arr.seek(0)

        # Add a slide to the presentation
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[5]
        )  # Use blank layout

        # Add the image to the slide
        left = top = Inches(0)
        slide.shapes.add_picture(
            img_byte_arr,
            left,
            top,
            width=Inches(pptx_width),
            height=Inches(pptx_height),
        )  # Adjust as needed

        if verbose >= 1:
            print(f"Slide {i}/{len(images)} added to the presentation")

    # Save the PowerPoint presentation
    presentation.save(output_pptx_path)

    if verbose >= 1:
        print(f"PDF has been converted to PPTX and saved at {output_pptx_path}")
